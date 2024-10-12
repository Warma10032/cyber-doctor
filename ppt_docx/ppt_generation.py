'''将大模型生成的json数据转换为ppt，可修改代码自定义ppt的样式'''
import datetime
import hashlib
import os
import time


from pptx.oxml.ns import qn
from typing import Dict
from pptx import Presentation
from env import get_app_root



_OUTPUT_DIR = os.path.join(get_app_root(), "data/cache/ppt")

# 如果文件夹路径不存在，先创建
if not os.path.exists(_OUTPUT_DIR):
    os.makedirs(_OUTPUT_DIR)

def get_file_path(text):
    """生成唯一的文件路径"""
    file_name = hashlib.sha256(text.encode("utf-8")).hexdigest()  ## 也可以使用uuid
    return os.path.join(_OUTPUT_DIR, f"{file_name}.pptx")

def generate(ppt_content: Dict) -> str:
    """生成 ppt 文件"""
    ppt = Presentation()

    # PPT首页
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout
    slide.placeholders[0].text = ppt_content["title"]
    slide.placeholders[1].text = "--来自「赛博华佗」"

    # 内容页
    print(f"总共{len(ppt_content['pages'])}页")
    for i, page in enumerate(ppt_content["pages"]):
        print("生成第%d页:%s" % (i + 1, page["title"]))
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout
        
        # 标题
        slide.placeholders[0].text = page["title"]
        # 正文
        text_frame = slide.placeholders[1].text_frame  # 获取文本框的text_frame对象
        

        for sub_content in page["content"]:
            print(sub_content)
            
            # 一级正文
            sub_title = text_frame.add_paragraph()
            sub_title.text, sub_title.level = sub_content["title"], 2
            
            # 二级正文
            sub_description = text_frame.add_paragraph()
            sub_description.text, sub_description.level = sub_content["description"], 3
            
    _output_file = get_file_path(str(time.time()))
    ppt.save(_output_file)

    return _output_file

